"""Raw visible report inventory for the verify skill.

HTML tables and visible blocks, Markdown lines, PDF lines, workbook cells,
and slide text use one stable item shape. Readers assign no claim meaning.
"""
from __future__ import annotations

import pathlib
import re
import sys
from html.parser import HTMLParser

SCRIPTS = pathlib.Path(__file__).resolve().parent
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from html_arith import _Tables  # noqa: E402

COMPLETED = frozenset({
    "confirmed", "contradicted", "not_checkable", "changed_since_report",
})
READERS = {
    ".html": "html",
    ".htm": "html",
    ".md": "md",
    ".txt": "txt",
    ".pdf": "pdf",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
}


def _visible_text(html: str) -> str:
    class _Text(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []
            self._skip = 0

        def handle_starttag(self, tag, attrs):
            if tag.lower() in {"script", "style"}:
                self._skip += 1

        def handle_endtag(self, tag):
            if tag.lower() in {"script", "style"} and self._skip:
                self._skip -= 1

        def handle_data(self, data):
            if not self._skip:
                self.parts.append(data)

    parser = _Text()
    parser.feed(html)
    parser.close()
    return re.sub(r"\s+", " ", "".join(parser.parts))


def _html_items(path: pathlib.Path) -> list[dict]:
    raw = path.read_text(errors="replace")
    items: list[dict] = []
    seen: set[tuple] = set()

    def add(kind: str, displayed: str, location: str) -> None:
        shown = re.sub(r"\s+", " ", displayed).strip()
        if not shown:
            return
        key = (kind, shown, location)
        if key in seen:
            return
        seen.add(key)
        items.append({
            "id": f"INV{len(items) + 1}",
            "kind": kind,
            "displayed": shown,
            "location": location,
            "importance": "unclassified",
            "quote": shown,
        })

    tables = _Tables()
    tables.feed(raw)
    tables.close()
    for t_i, table in enumerate(tables.tables, start=1):
        for r_i, row in enumerate(table):
            for c_i, cell in enumerate(row):
                add("table_cell", cell, f"table{t_i}/r{r_i + 1}/c{c_i + 1}")

    class _Occurrences(HTMLParser):
        """Collect each visible body text occurrence without assigning meaning."""

        def __init__(self) -> None:
            super().__init__()
            self.occurrences: list[str] = []
            self._hidden = 0
            self._table = 0

        def handle_starttag(self, tag, attrs):
            tag = tag.lower()
            if tag in {"head", "script", "style", "template"}:
                self._hidden += 1
            elif tag == "table":
                self._table += 1

        def handle_endtag(self, tag):
            tag = tag.lower()
            if tag in {"head", "script", "style", "template"} and self._hidden:
                self._hidden -= 1
            elif tag == "table" and self._table:
                self._table -= 1

        def handle_data(self, data):
            if self._hidden or self._table:
                return
            shown = re.sub(r"\s+", " ", data).strip()
            if shown:
                self.occurrences.append(shown)

    occurrences = _Occurrences()
    occurrences.feed(raw)
    occurrences.close()
    for index, shown in enumerate(occurrences.occurrences, 1):
        add("html_text", shown, f"text{index}")
    return items


def _bag_add(items: list, seen: set, kind: str, displayed: str,
             location: str) -> None:
    shown = re.sub(r"\s+", " ", displayed).strip()
    if not shown:
        return
    key = (kind, shown, location)
    if key in seen:
        return
    seen.add(key)
    items.append({
        "id": f"INV{len(items) + 1}",
        "kind": kind,
        "displayed": shown,
        "location": location,
        "importance": "unclassified",
        "quote": shown,
    })


def _md_items(path: pathlib.Path) -> list[dict]:
    items: list[dict] = []
    seen: set[tuple] = set()
    for index, raw in enumerate(path.read_text(errors="replace").splitlines(), 1):
        line = raw.strip()
        if not line:
            continue
        shown = re.sub(r"^(?:#{1,6}|[-*])\s+", "", line).strip()
        if not shown:
            continue
        _bag_add(items, seen, "md_line", shown, f"line{index}")
    return items


def visible_markdown(path: pathlib.Path) -> str:
    return path.read_text(errors="replace")


def _pdf_items(path: pathlib.Path) -> tuple[list[dict], str, str | None]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return [], "", "pypdf is not installed"
    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001 — fail closed on unreadable PDF
        return [], "", f"unreadable PDF: {exc}"
    items: list[dict] = []
    seen: set[tuple] = set()
    pages: list[str] = []
    for page_i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append(text)
        for line_i, raw in enumerate(text.splitlines(), start=1):
            shown = re.sub(r"\s+", " ", raw).strip()
            if not shown:
                continue
            _bag_add(
                items, seen, "pdf_line", shown, f"page{page_i}/line{line_i}")
    visible = "\n".join(pages).strip()
    if not visible:
        return [], "", "no extractable PDF text"
    return items, visible, None


def _xlsx_display(cell) -> str | None:
    value = cell.value
    if value is None or value == "":
        return None
    fmt = str(cell.number_format or "General")
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, str):
        shown = re.sub(r"\s+", " ", value).strip()
        return shown or None
    if isinstance(value, (int, float)):
        if "%" in fmt:
            match = re.search(r"0\.(0+)%", fmt)
            places = len(match.group(1)) if match else 0
            return f"{float(value) * 100:.{places}f}%"
        if "0.00" in fmt:
            return f"{float(value):,.2f}"
        if "#,##0" in fmt:
            return f"{float(value):,.0f}"
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)
    return str(value)


def _xlsx_items(path: pathlib.Path) -> tuple[list[dict], str, str | None]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return [], "", "openpyxl is not installed"
    try:
        book = load_workbook(path, data_only=False)
    except Exception as exc:  # noqa: BLE001
        return [], "", f"unreadable xlsx: {exc}"
    items: list[dict] = []
    seen: set[tuple] = set()
    lines: list[str] = []
    for sheet in book.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                shown = _xlsx_display(cell)
                if shown is None:
                    continue
                lines.append(f"{sheet.title}!{cell.coordinate} {shown}")
                _bag_add(
                    items, seen, "xlsx_cell", shown,
                    f"{sheet.title}/{cell.coordinate}")
    visible = "\n".join(lines)
    if not visible:
        return [], "", "no visible xlsx cells"
    return items, visible, None


def _pptx_items(path: pathlib.Path) -> tuple[list[dict], str, str | None]:
    try:
        from pptx import Presentation
    except ImportError:
        return [], "", "python-pptx is not installed"
    try:
        deck = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return [], "", f"unreadable pptx: {exc}"
    items: list[dict] = []
    seen: set[tuple] = set()
    lines: list[str] = []
    for slide_i, slide in enumerate(deck.slides, start=1):
        for shape_i, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = re.sub(r"\s+", " ", shape.text_frame.text or "").strip()
            if not text:
                continue
            lines.append(text)
            _bag_add(
                items, seen, "pptx_shape", text,
                f"slide{slide_i}/shape{shape_i}")
    visible = "\n".join(lines)
    if not visible:
        return [], "", "no visible pptx text"
    return items, visible, None


def visible_text_for(path: pathlib.Path) -> tuple[str, str | None]:
    """Visible report text plus an error when extraction failed."""
    suffix = path.suffix.lower()
    reader = READERS.get(suffix, suffix.lstrip(".") or "unknown")
    if reader == "html":
        return _visible_text(path.read_text(errors="replace")), None
    if reader in {"md", "txt"}:
        return visible_markdown(path), None
    if reader == "pdf":
        _items, visible, err = _pdf_items(path)
        return visible, err
    if reader == "xlsx":
        _items, visible, err = _xlsx_items(path)
        return visible, err
    if reader == "pptx":
        _items, visible, err = _pptx_items(path)
        return visible, err
    return "", f"no visible-text reader for {reader}"


def inventory_for(path: pathlib.Path) -> dict:
    """Same shape for every supported format."""
    suffix = path.suffix.lower()
    reader = READERS.get(suffix, suffix.lstrip(".") or "unknown")
    if reader == "html":
        return {
            "reader": "html",
            "complete": True,
            "items": _html_items(path),
            "reason": None,
        }
    if reader in {"md", "txt"}:
        items = _md_items(path)
        return {
            "reader": reader,
            "complete": True,
            "items": items,
            "reason": None,
        }
    if reader == "pdf":
        items, _visible, err = _pdf_items(path)
        return {
            "reader": "pdf",
            "complete": err is None,
            "items": items if err is None else [],
            "reason": err,
        }
    if reader == "xlsx":
        items, _visible, err = _xlsx_items(path)
        return {
            "reader": "xlsx",
            "complete": err is None,
            "items": items if err is None else [],
            "reason": err,
        }
    if reader == "pptx":
        items, _visible, err = _pptx_items(path)
        return {
            "reader": "pptx",
            "complete": err is None,
            "items": items if err is None else [],
            "reason": err,
        }
    return {
        "reader": reader,
        "complete": False,
        "items": [],
        "reason": f"no inventory reader for {reader}",
    }


def claim_inventory_ids(claim: dict) -> list[str]:
    raw = claim.get("inventory_ids")
    if isinstance(raw, str):
        return [raw.strip()] if raw.strip() else []
    if isinstance(raw, list):
        out = []
        seen = set()
        for value in raw:
            item = str(value or "").strip()
            if item and item not in seen:
                seen.add(item)
                out.append(item)
        return out
    return []


def cover(inventory: dict, claims: list, *,
          structural_context: list[dict] | None = None) -> dict:
    """Reconcile explicit host classifications by inventory id only."""
    items = [
        item for item in (inventory.get("items") or [])
        if isinstance(item, dict)
    ]
    by_id = {str(item.get("id") or ""): item for item in items if item.get("id")}
    consumed: dict[str, dict] = {}
    missing: list[dict] = []
    mapping: list[dict] = []
    assignments = list(claims) + list(structural_context or [])
    for assignment in assignments:
        assignment_id = str(
            assignment.get("id")
            or assignment.get("candidate_id")
            or "structural_context"
        )
        classification = str(assignment.get("classification") or "")
        outcome = assignment.get("outcome")
        for iid in claim_inventory_ids(assignment):
            item = by_id.get(iid)
            if item is None:
                missing.append({
                    "id": iid,
                    "claim_id": assignment_id,
                    "reason": "inventory id is not present",
                })
                continue
            if iid in consumed:
                missing.append({
                    "id": iid,
                    "displayed": item.get("displayed"),
                    "location": item.get("location"),
                    "claim_id": assignment_id,
                    "reason": "inventory id is consumed more than once",
                })
                continue
            consumed[iid] = assignment
            if (
                item.get("importance") == "unclassified"
                or item.get("classification") != classification
            ):
                missing.append({
                    "id": iid,
                    "displayed": item.get("displayed"),
                    "location": item.get("location"),
                    "claim_id": assignment_id,
                    "reason": "inventory classification was not applied exactly",
                })
            mapping.append({
                "inventory_id": iid,
                "claim_id": assignment_id,
                "classification": classification,
                "outcome": outcome,
            })
    for item in items:
        iid = str(item.get("id") or "")
        if iid and iid not in consumed:
            missing.append({
                "id": iid,
                "displayed": item.get("displayed"),
                "location": item.get("location"),
                "reason": "inventory item has no host classification",
            })
    material_items = [
        item for item in items
        if item.get("classification") == "material_claim"
        and item.get("importance") == "material"
    ]
    material_ids = {str(item.get("id") or "") for item in material_items}
    accounted = sum(iid in consumed for iid in material_ids)
    completed = sum(
        iid in consumed and consumed[iid].get("outcome") in COMPLETED
        for iid in material_ids
    )
    for iid in material_ids:
        assignment = consumed.get(iid)
        if assignment is not None and assignment.get("outcome") not in COMPLETED:
            item = by_id[iid]
            missing.append({
                "id": iid,
                "displayed": item.get("displayed"),
                "location": item.get("location"),
                "claim_id": assignment.get("id"),
                "outcome": assignment.get("outcome"),
                "reason": "material inventory item has no completed outcome",
            })
    n = len(material_items)
    return {
        "material": n,
        "accounted": accounted,
        "completed": completed,
        "missing": missing,
        "mapping": mapping,
        "extractor_fraction": (accounted / n) if n else 1.0,
        "engine_fraction": (completed / n) if n else 1.0,
        "inventory_complete": bool(inventory.get("complete")),
    }
